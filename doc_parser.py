import os
import logging
from pypdf import PdfReader

# Set up logging
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)


def chunk_text(text: str, chunk_size_words: int = 350, overlap_words: int = 70) -> list[str]:
    """Chunks text by words with a specified overlap."""
    words = text.split()
    chunks = []
    
    if len(words) <= chunk_size_words:
        return [" ".join(words)]
        
    start = 0
    while start < len(words):
        end = start + chunk_size_words
        chunk = words[start:end]
        chunks.append(" ".join(chunk))
        
        if end >= len(words):
            break
        start += (chunk_size_words - overlap_words)
        
    return chunks


def parse_pdf(file_path: str) -> list[dict]:
    """Parses PDF page by page."""
    chunks_data = []
    filename = os.path.basename(file_path)
    
    try:
        reader = PdfReader(file_path)
        for page_idx, page in enumerate(reader.pages):
            text = page.extract_text()
            if not text or not text.strip():
                continue
                
            clean_text = " ".join(text.split())
            page_chunks = chunk_text(clean_text)
            
            for chunk_idx, chunk in enumerate(page_chunks):
                chunks_data.append({
                    "text": chunk,
                    "metadata": {
                        "source_file": filename,
                        "page_number": page_idx + 1,
                        "chunk_index": chunk_idx
                    }
                })
    except Exception as e:
        logger.error(f"Error parsing PDF '{file_path}': {e}")
        
    return chunks_data


def parse_docx(file_path: str) -> list[dict]:
    """Parses DOCX page by page (approximated by grouping text into ~350-word pages)."""
    chunks_data = []
    filename = os.path.basename(file_path)
    
    try:
        import docx
        doc = docx.Document(file_path)
        
        # Gather all text from paragraphs and tables
        full_text_list = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text_list.append(para.text.strip())
                
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        full_text_list.append(cell.text.strip())
                        
        full_text = " ".join(full_text_list)
        words = full_text.split()
        
        # Approximate pages by grouping every 350 words as a "page"
        words_per_page = 350
        total_pages = max(1, len(words) // words_per_page + (1 if len(words) % words_per_page > 0 else 0))
        
        for page_idx in range(total_pages):
            start_word = page_idx * words_per_page
            end_word = start_word + words_per_page
            page_words = words[start_word:end_word]
            page_text = " ".join(page_words)
            
            if not page_text.strip():
                continue
                
            page_chunks = chunk_text(page_text)
            for chunk_idx, chunk in enumerate(page_chunks):
                chunks_data.append({
                    "text": chunk,
                    "metadata": {
                        "source_file": filename,
                        "page_number": page_idx + 1,
                        "chunk_index": chunk_idx
                    }
                })
    except ImportError:
        logger.error("python-docx not installed. Cannot parse DOCX files.")
    except Exception as e:
        logger.error(f"Error parsing DOCX '{file_path}': {e}")
        
    return chunks_data


def parse_pptx(file_path: str) -> list[dict]:
    """Parses PPTX slides (each slide counts as a page/slide number)."""
    chunks_data = []
    filename = os.path.basename(file_path)
    
    try:
        import pptx
        prs = pptx.Presentation(file_path)
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text_list = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                slide_text_list.append(run.text.strip())
            
            slide_text = " ".join(slide_text_list)
            if not slide_text.strip():
                continue
                
            slide_chunks = chunk_text(slide_text)
            for chunk_idx, chunk in enumerate(slide_chunks):
                chunks_data.append({
                    "text": chunk,
                    "metadata": {
                        "source_file": filename,
                        "page_number": slide_idx + 1,  # Slide number
                        "chunk_index": chunk_idx
                    }
                })
    except ImportError:
        logger.error("python-pptx not installed. Cannot parse PPTX files.")
    except Exception as e:
        logger.error(f"Error parsing PPTX '{file_path}': {e}")
        
    return chunks_data


def parse_document(file_path: str) -> list[dict]:
    """General function to parse document depending on extension."""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext in [".docx", ".doc"]:
        return parse_docx(file_path)
    elif ext in [".pptx", ".ppt"]:
        return parse_pptx(file_path)
    else:
        logger.warning(f"Unsupported file format '{ext}' for file: '{file_path}'")
        return []


if __name__ == "__main__":
    # Test script locally
    test_pdf = os.path.join("data", "ethicsds.pdf")
    if os.path.exists(test_pdf):
        print(f"Testing document parser on: {test_pdf}")
        res = parse_document(test_pdf)
        print(f"Parsed {len(res)} chunks successfully.")
        if res:
            print("First chunk sample metadata:", res[0]["metadata"])
            print("First chunk text sample:", res[0]["text"][:150] + "...")
